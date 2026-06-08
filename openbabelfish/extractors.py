"""
Document extraction module for OpenBabelFish.

Converts PDF, DOCX, PPTX, and EPUB files into plain text for translation.
Supports both text-based and scanned/image PDFs via EasyOCR fallback.
"""

import logging
from pathlib import Path
from typing import Optional

from rich.console import Console
from rich.prompt import Confirm

logger = logging.getLogger(__name__)
console = Console()

# Minimum average characters per page to consider a PDF as text-based.
# Below this threshold, we assume the PDF is image-based and needs OCR.
TEXT_PDF_CHAR_THRESHOLD = 50

# Supported file extensions and their extractor mapping
SUPPORTED_EXTENSIONS = {
    ".txt": "text",
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".epub": "epub",
}

# Dependencies required for each extraction type
EXTRACTION_DEPS = {
    "pdf": ["PyMuPDF"],
    "docx": ["python-docx"],
    "pptx": ["python-pptx"],
    "epub": ["EbookLib", "beautifulsoup4"],
    "ocr": ["easyocr"],
}


class ExtractionError(Exception):
    """Raised when text extraction from a document fails."""


class FileExtractor:
    """
    Top-level router that detects file type by extension and delegates
    to the appropriate extractor.

    Handles dependency checking and auto-install prompting before
    extraction begins.
    """

    def __init__(self, dep_mgr=None):
        """
        Args:
            dep_mgr: A DependencyManager instance for auto-installing
                     missing extraction packages. If None, missing
                     packages will raise an error instead of prompting.
        """
        self._dep_mgr = dep_mgr

    def extract(self, file_path: str, force_ocr: bool = False) -> str:
        """
        Extract plain text from a document file.

        Args:
            file_path: Path to the input file.
            force_ocr: If True and the file is a PDF, skip text extraction
                       and go straight to OCR.

        Returns:
            The extracted plain text.

        Raises:
            FileNotFoundError: If the file does not exist.
            ExtractionError: If extraction fails or the format is unsupported.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            supported = ", ".join(SUPPORTED_EXTENSIONS.keys())
            raise ExtractionError(
                f"Unsupported file format: '{ext}'. "
                f"Supported formats: {supported}"
            )

        kind = SUPPORTED_EXTENSIONS[ext]

        if kind == "text":
            return path.read_text(encoding="utf-8")

        # Ensure the required packages are installed before extraction
        self._ensure_deps(kind)
        if force_ocr and kind == "pdf":
            self._ensure_deps("ocr")

        if kind == "pdf":
            return PDFExtractor.extract(str(path), force_ocr=force_ocr)
        elif kind == "docx":
            return DOCXExtractor.extract(str(path))
        elif kind == "pptx":
            return PPTXExtractor.extract(str(path))
        elif kind == "epub":
            return EPUBExtractor.extract(str(path))

        raise ExtractionError(f"No extractor available for: {ext}")

    def _ensure_deps(self, kind: str) -> None:
        """Check and auto-install dependencies for a given extraction kind."""
        deps = EXTRACTION_DEPS.get(kind, [])
        if not deps:
            return

        missing = []
        for pkg in deps:
            if not _is_package_installed(pkg):
                missing.append(pkg)

        if not missing:
            return

        if self._dep_mgr is None:
            raise ExtractionError(
                f"Missing packages for {kind} extraction: {', '.join(missing)}. "
                f"Install them with: pip install {' '.join(missing)}"
            )

        pkg_list = ", ".join(missing)
        console.print(
            f"\n[yellow]⚠  Document extraction requires: [cyan]{pkg_list}[/cyan][/yellow]"
        )
        if Confirm.ask("  Install now?", default=True):
            self._dep_mgr.install_missing(missing)
        else:
            raise ExtractionError(
                f"Cannot extract without packages: {pkg_list}"
            )


class PDFExtractor:
    """
    Extracts text from PDF files.

    Strategy:
    1. Attempt text extraction via PyMuPDF (fast, preserves layout).
    2. If the PDF appears to be image-based (low text content), fall back
       to EasyOCR.
    3. The --ocr flag bypasses step 1 and goes directly to OCR.
    """

    @staticmethod
    def extract(file_path: str, force_ocr: bool = False) -> str:
        """
        Extract text from a PDF.

        Args:
            file_path: Path to the PDF file.
            force_ocr: If True, skip text extraction and use OCR directly.

        Returns:
            Extracted plain text from all pages.
        """
        if force_ocr:
            console.print("[dim]  OCR mode forced by --ocr flag.[/dim]")
            return PDFExtractor._extract_ocr(file_path)

        # Try text extraction first
        text = PDFExtractor._extract_text(file_path)

        if PDFExtractor._is_text_pdf(file_path, text):
            return text

        # Text extraction yielded too little content — fall back to OCR
        console.print(
            "[yellow]  PDF appears to be image-based. "
            "Falling back to OCR...[/yellow]"
        )
        # Ensure OCR deps are available before attempting
        if not _is_package_installed("easyocr"):
            console.print(
                "\n[yellow]⚠  OCR requires EasyOCR to be installed.[/yellow]"
            )
            if Confirm.ask("  Install EasyOCR now?", default=True):
                from .managers import DependencyManager
                DependencyManager.install_ocr_support()
            else:
                # Return whatever text we got, even if sparse
                console.print(
                    "[dim]  Returning partial text extraction.[/dim]"
                )
                return text

        return PDFExtractor._extract_ocr(file_path)

    @staticmethod
    def _extract_text(file_path: str) -> str:
        """Extract selectable text from PDF pages using PyMuPDF."""
        import fitz  # PyMuPDF

        pages = []
        with fitz.open(file_path) as doc:
            for page in doc:
                page_text = page.get_text("text")
                if page_text.strip():
                    pages.append(page_text.strip())

        return "\n\n".join(pages)

    @staticmethod
    def _extract_ocr(file_path: str) -> str:
        """Extract text from PDF pages as images using EasyOCR."""
        import fitz  # PyMuPDF
        import easyocr
        from .config import load_config

        config = load_config()
        ocr_device = config.get("ocr_device", "cpu")

        gpu_enabled = False
        if ocr_device == "gpu":
            try:
                import torch
                if torch.cuda.is_available():
                    gpu_enabled = True
                else:
                    console.print("[yellow]  ⚠  CUDA is not available for PyTorch. Falling back to CPU for OCR.[/yellow]")
            except ImportError:
                console.print("[yellow]  ⚠  PyTorch is not installed. Falling back to CPU for OCR.[/yellow]")

        reader = easyocr.Reader(["en"], gpu=gpu_enabled)
        pages = []

        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                # Render page to a high-res image (300 DPI)
                mat = fitz.Matrix(300 / 72, 300 / 72)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")

                # Run OCR on the image bytes
                results = reader.readtext(img_bytes, detail=0)
                page_text = " ".join(results)

                if page_text.strip():
                    pages.append(page_text.strip())

                logger.debug(
                    "OCR page %d: extracted %d characters",
                    page_num + 1,
                    len(page_text),
                )

        return "\n\n".join(pages)

    @staticmethod
    def _is_text_pdf(file_path: str, extracted_text: Optional[str] = None) -> bool:
        """
        Heuristic to determine if a PDF has meaningful selectable text.

        Checks whether the average character count per page exceeds
        TEXT_PDF_CHAR_THRESHOLD. If it does, the PDF is considered
        text-based; otherwise, it's likely a scanned document.
        """
        import fitz  # PyMuPDF

        if extracted_text is None:
            extracted_text = PDFExtractor._extract_text(file_path)

        with fitz.open(file_path) as doc:
            num_pages = len(doc)

        if num_pages == 0:
            return False

        avg_chars = len(extracted_text) / num_pages
        return avg_chars >= TEXT_PDF_CHAR_THRESHOLD


class DOCXExtractor:
    """Extracts text from Microsoft Word (.docx) files."""

    @staticmethod
    def extract(file_path: str) -> str:
        """
        Extract text from a DOCX file.

        Reads all paragraphs in document order, preserving paragraph
        breaks as double newlines.
        """
        from docx import Document

        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        return "\n\n".join(paragraphs)


class PPTXExtractor:
    """Extracts text from PowerPoint (.pptx) files."""

    @staticmethod
    def extract(file_path: str) -> str:
        """
        Extract text from a PPTX file.

        Iterates through all slides and text frames, collecting text
        from each shape. Slides are separated by double newlines.
        """
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)

            if slide_parts:
                slides_text.append("\n".join(slide_parts))

        return "\n\n".join(slides_text)


class EPUBExtractor:
    """Extracts text from EPUB e-book files."""

    @staticmethod
    def extract(file_path: str) -> str:
        """
        Extract text from an EPUB file.

        Reads all document items (XHTML chapters) and strips HTML tags
        using BeautifulSoup, preserving paragraph structure.
        """
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup

        book = epub.read_epub(file_path, options={"ignore_ncx": True})
        chapters = []

        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            content = item.get_content().decode("utf-8", errors="replace")
            soup = BeautifulSoup(content, "html.parser")

            # Extract text from paragraph and heading tags for structure
            parts = []
            for tag in soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6"]):
                text = tag.get_text(strip=True)
                if text:
                    parts.append(text)

            if parts:
                chapters.append("\n\n".join(parts))

        return "\n\n".join(chapters)


def _is_package_installed(package_name: str) -> bool:
    """Check if a Python package is importable."""
    import importlib.metadata

    # Map pip package names to their importable names
    _import_map = {
        "PyMuPDF": "fitz",
        "python-docx": "docx",
        "python-pptx": "pptx",
        "EbookLib": "ebooklib",
        "beautifulsoup4": "bs4",
        "easyocr": "easyocr",
    }

    # Try import-based check first (most reliable)
    import_name = _import_map.get(package_name, package_name)
    try:
        __import__(import_name)
        return True
    except ImportError:
        pass

    # Fallback: try metadata lookup with the pip package name
    try:
        importlib.metadata.version(package_name)
        return True
    except importlib.metadata.PackageNotFoundError:
        return False
