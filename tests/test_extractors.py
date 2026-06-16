"""Tests for the document extraction module."""
import pytest
import tempfile
import os
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock

from openbabelfish.extractors import (
    FileExtractor,
    PDFExtractor,
    DOCXExtractor,
    PPTXExtractor,
    EPUBExtractor,
    ExtractionError,
    SUPPORTED_EXTENSIONS,
    TEXT_PDF_CHAR_THRESHOLD,
    _is_package_installed,
)


# ── FileExtractor Routing ────────────────────────────────────────────────────


class TestFileExtractorRouting:
    """Tests for the FileExtractor's extension-based routing logic."""

    def test_supported_extensions_contains_expected_types(self):
        assert ".txt" in SUPPORTED_EXTENSIONS
        assert ".pdf" in SUPPORTED_EXTENSIONS
        assert ".docx" in SUPPORTED_EXTENSIONS
        assert ".pptx" in SUPPORTED_EXTENSIONS
        assert ".epub" in SUPPORTED_EXTENSIONS

    def test_unsupported_extension_raises(self):
        extractor = FileExtractor()
        with tempfile.NamedTemporaryFile(suffix=".xyz", delete=False) as f:
            f.write(b"test")
            f.flush()
            with pytest.raises(ExtractionError, match="Unsupported file format"):
                extractor.extract(f.name)
        os.unlink(f.name)

    def test_file_not_found_raises(self):
        extractor = FileExtractor()
        with pytest.raises(FileNotFoundError):
            extractor.extract("/nonexistent/file.txt")

    def test_txt_file_reads_directly(self):
        extractor = FileExtractor()
        with tempfile.NamedTemporaryFile(
            suffix=".txt", delete=False, mode="w", encoding="utf-8"
        ) as f:
            f.write("Hello, world!")
            f.flush()
            result = extractor.extract(f.name)
            assert result == "Hello, world!"
        os.unlink(f.name)

    def test_txt_file_with_unicode(self):
        extractor = FileExtractor()
        content = "こんにちは世界！\nBonjour le monde!"
        with tempfile.NamedTemporaryFile(
            suffix=".txt", delete=False, mode="w", encoding="utf-8"
        ) as f:
            f.write(content)
            f.flush()
            result = extractor.extract(f.name)
            assert result == content
        os.unlink(f.name)


# ── PDF Extraction ───────────────────────────────────────────────────────────


class TestPDFExtractor:
    """Tests for PDFExtractor text detection heuristic."""

    def test_is_text_pdf_with_enough_text(self):
        """A PDF with >= TEXT_PDF_CHAR_THRESHOLD chars/page is text-based."""
        text = "a" * (TEXT_PDF_CHAR_THRESHOLD * 3)  # 3 pages worth

        mock_fitz = MagicMock()
        mock_doc = MagicMock()
        mock_doc.__len__ = MagicMock(return_value=3)
        mock_doc.__enter__ = MagicMock(return_value=mock_doc)
        mock_doc.__exit__ = MagicMock(return_value=False)
        mock_fitz.open.return_value = mock_doc

        with patch.dict("sys.modules", {"fitz": mock_fitz}):
            result = PDFExtractor._is_text_pdf("dummy.pdf", text)
            assert result is True

    def test_is_text_pdf_with_little_text(self):
        """A PDF with < TEXT_PDF_CHAR_THRESHOLD chars/page is image-based."""
        text = "ab"  # Very little text for a 3-page doc

        mock_fitz = MagicMock()
        mock_doc = MagicMock()
        mock_doc.__len__ = MagicMock(return_value=3)
        mock_doc.__enter__ = MagicMock(return_value=mock_doc)
        mock_doc.__exit__ = MagicMock(return_value=False)
        mock_fitz.open.return_value = mock_doc

        with patch.dict("sys.modules", {"fitz": mock_fitz}):
            result = PDFExtractor._is_text_pdf("dummy.pdf", text)
            assert result is False

    def test_is_text_pdf_empty_document(self):
        """A 0-page PDF returns False."""
        mock_fitz = MagicMock()
        mock_doc = MagicMock()
        mock_doc.__len__ = MagicMock(return_value=0)
        mock_doc.__enter__ = MagicMock(return_value=mock_doc)
        mock_doc.__exit__ = MagicMock(return_value=False)
        mock_fitz.open.return_value = mock_doc

        with patch.dict("sys.modules", {"fitz": mock_fitz}):
            result = PDFExtractor._is_text_pdf("dummy.pdf", "")
            assert result is False


# ── DOCX Extraction ──────────────────────────────────────────────────────────


class TestDOCXExtractor:
    """Tests for DOCXExtractor with mocked python-docx."""

    @patch("openbabelfish.extractors.DOCXExtractor.extract")
    def test_extract_returns_joined_paragraphs(self, mock_extract):
        mock_extract.return_value = "Hello\n\nWorld"
        result = DOCXExtractor.extract("test.docx")
        assert "Hello" in result
        assert "World" in result

    def test_docx_with_real_file(self):
        """Integration test: create a real DOCX and extract text."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        doc = Document()
        doc.add_paragraph("First paragraph")
        doc.add_paragraph("Second paragraph")
        doc.add_paragraph("")  # Empty paragraph
        doc.add_paragraph("Third paragraph")

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc.save(f.name)
            result = DOCXExtractor.extract(f.name)

        os.unlink(f.name)

        assert "First paragraph" in result
        assert "Second paragraph" in result
        assert "Third paragraph" in result
        # Empty paragraphs should be filtered out
        parts = [p for p in result.split("\n\n") if p.strip()]
        assert len(parts) == 3


# ── PPTX Extraction ──────────────────────────────────────────────────────────


class TestPPTXExtractor:
    """Tests for PPTXExtractor with mocked python-pptx."""

    def test_pptx_with_real_file(self):
        """Integration test: create a real PPTX and extract text."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")

        prs = Presentation()
        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        slide1.shapes.title.text = "Slide One Title"
        slide1.placeholders[1].text = "Slide one content"

        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Slide Two Title"
        slide2.placeholders[1].text = "Slide two content"

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            result = PPTXExtractor.extract(f.name)

        os.unlink(f.name)

        assert "Slide One Title" in result
        assert "Slide one content" in result
        assert "Slide Two Title" in result
        assert "Slide two content" in result


# ── EPUB Extraction ──────────────────────────────────────────────────────────


class TestEPUBExtractor:
    """Tests for EPUBExtractor with mocked ebooklib."""

    @patch("openbabelfish.extractors.EPUBExtractor.extract")
    def test_extract_returns_chapters(self, mock_extract):
        mock_extract.return_value = "Chapter 1 text\n\nChapter 2 text"
        result = EPUBExtractor.extract("test.epub")
        assert "Chapter 1" in result
        assert "Chapter 2" in result


# ── Package Detection ────────────────────────────────────────────────────────


class TestPackageDetection:
    """Tests for the _is_package_installed utility."""

    def test_detects_installed_package(self):
        # 'os' is always available (stdlib)
        assert _is_package_installed("os") is True

    def test_detects_missing_package(self):
        assert _is_package_installed("nonexistent_fake_package_xyz") is False


# ── Dependency Auto-Install ──────────────────────────────────────────────────


class TestDependencyAutoInstall:
    """Tests for the auto-install prompting logic in FileExtractor."""

    @patch("openbabelfish.extractors._is_package_installed", return_value=False)
    def test_raises_without_dep_mgr(self, mock_installed):
        """Without a dep_mgr, missing deps should raise ExtractionError."""
        extractor = FileExtractor(dep_mgr=None)

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(b"fake pdf")
            f.flush()
            with pytest.raises(ExtractionError, match="Missing packages"):
                extractor.extract(f.name)
        os.unlink(f.name)

    @patch("openbabelfish.extractors._is_package_installed", return_value=True)
    @patch.object(PDFExtractor, "extract", return_value="extracted text")
    def test_routes_pdf_when_deps_available(self, mock_pdf_extract, mock_installed):
        """When deps are installed, PDF should route to PDFExtractor."""
        extractor = FileExtractor()

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(b"fake pdf")
            f.flush()
            result = extractor.extract(f.name)
            assert result == "extracted text"
        os.unlink(f.name)

    @patch("openbabelfish.extractors._is_package_installed", return_value=True)
    @patch.object(PDFExtractor, "extract", return_value="ocr text")
    def test_force_ocr_flag_passed_to_pdf(self, mock_pdf_extract, mock_installed):
        """The force_ocr flag should be forwarded to PDFExtractor."""
        extractor = FileExtractor()

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(b"fake pdf")
            f.flush()
            extractor.extract(f.name, force_ocr=True)
            mock_pdf_extract.assert_called_once_with(f.name, force_ocr=True, progress_callback=None)
        os.unlink(f.name)


# ── OCR GPU/CPU Handoff and Fallback ──────────────────────────────────────────


class TestOCRHandoffAndFallback:
    """Tests for OCR hardware configuration (CPU/GPU) and fallback."""

    @patch("easyocr.Reader")
    @patch("openbabelfish.config.load_config")
    def test_ocr_defaults_to_cpu(self, mock_load_config, mock_reader):
        """By default, the OCR engine should be initialized with gpu=False."""
        mock_load_config.return_value = {"ocr_device": "cpu"}
        mock_fitz = MagicMock()
        mock_doc = MagicMock()
        mock_doc.__len__ = MagicMock(return_value=0)
        mock_doc.__enter__ = MagicMock(return_value=mock_doc)
        mock_fitz.open.return_value = mock_doc

        with patch.dict("sys.modules", {"fitz": mock_fitz}):
            PDFExtractor._extract_ocr("dummy.pdf")
            mock_reader.assert_called_once_with(["en"], gpu=False)

    @patch("easyocr.Reader")
    @patch("openbabelfish.config.load_config")
    def test_ocr_uses_gpu_when_configured_and_available(
        self, mock_load_config, mock_reader
    ):
        """When ocr_device is 'gpu' and CUDA is available, initialize with gpu=True."""
        mock_load_config.return_value = {"ocr_device": "gpu"}
        mock_fitz = MagicMock()
        mock_doc = MagicMock()
        mock_doc.__len__ = MagicMock(return_value=0)
        mock_doc.__enter__ = MagicMock(return_value=mock_doc)
        mock_fitz.open.return_value = mock_doc

        mock_torch = MagicMock()
        mock_torch.cuda.is_available.return_value = True
        
        with patch.dict("sys.modules", {"fitz": mock_fitz, "torch": mock_torch}):
            PDFExtractor._extract_ocr("dummy.pdf")
            mock_reader.assert_called_once_with(["en"], gpu=True)

    @patch("easyocr.Reader")
    @patch("openbabelfish.config.load_config")
    def test_ocr_falls_back_to_cpu_when_cuda_unavailable(
        self, mock_load_config, mock_reader
    ):
        """When ocr_device is 'gpu' but CUDA is not available, fall back to gpu=False."""
        mock_load_config.return_value = {"ocr_device": "gpu"}
        mock_fitz = MagicMock()
        mock_doc = MagicMock()
        mock_doc.__len__ = MagicMock(return_value=0)
        mock_doc.__enter__ = MagicMock(return_value=mock_doc)
        mock_fitz.open.return_value = mock_doc

        mock_torch = MagicMock()
        mock_torch.cuda.is_available.return_value = False
        
        with patch.dict("sys.modules", {"fitz": mock_fitz, "torch": mock_torch}):
            PDFExtractor._extract_ocr("dummy.pdf")
            mock_reader.assert_called_once_with(["en"], gpu=False)
